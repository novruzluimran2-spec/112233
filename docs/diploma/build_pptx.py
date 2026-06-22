"""Generate PowerPoint presentation for diploma defense."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "presentation-antalya.pptx"

BG = RGBColor(0x0F, 0x17, 0x2A)
GOLD = RGBColor(0xFB, 0xBF, 0x24)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x86, 0xEF, 0xAC)
CARD = RGBColor(0x1E, 0x29, 0x3B)


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, items, size=16, color=WHITE):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tf


def add_title_slide(prs, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    y = Inches(1.2)
    for i, (text, size, bold, color, align) in enumerate(lines):
        add_textbox(slide, Inches(0.8), y, Inches(8.4), Inches(1.2), text, size, bold, color, align)
        y += Inches(0.55 if size <= 20 else 0.85)


def add_content_slide(prs, title, bullets=None, paragraphs=None, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(9), Inches(0.7), title, 28, True, GOLD)
    y = Inches(1.15)
    if paragraphs:
        for para in paragraphs:
            add_textbox(slide, Inches(0.8), y, Inches(8.6), Inches(1.5), para, 17, False, WHITE)
            y += Inches(0.55)
    if bullets:
        add_bullets(slide, Inches(0.8), y, Inches(8.6), Inches(4.5), bullets, 17)
    if footer:
        add_textbox(slide, Inches(0.8), Inches(6.5), Inches(8.6), Inches(0.6), footer, 15, True, GREEN)


def add_table_slide(prs, title, headers, rows, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(9), Inches(0.7), title, 28, True, GOLD)
    cols, row_count = len(headers), len(rows) + 1
    table = slide.shapes.add_table(row_count, cols, Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.45 * row_count)).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = GOLD
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = WHITE
    if note:
        add_textbox(slide, Inches(0.8), Inches(5.8), Inches(8.6), Inches(0.8), note, 15, False, MUTED)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, [
        ("ГАПОУ «Волгоградский социально-педагогический колледж»", 16, False, MUTED, PP_ALIGN.CENTER),
        ("Специальность 09.02.07 — Информационные системы и программирование", 14, False, MUTED, PP_ALIGN.CENTER),
        ("Разработка сайта для ресторана «Анталия»", 32, True, GOLD, PP_ALIGN.CENTER),
        ("Выполнил: студент группы 43 Дд\nНоврузлу Имран Мехман Оглы", 18, False, WHITE, PP_ALIGN.CENTER),
        ("Руководитель: Авдосиева Светалана Васильевна", 16, False, WHITE, PP_ALIGN.CENTER),
        ("Волгоград, 2026", 14, False, MUTED, PP_ALIGN.CENTER),
    ])

    add_content_slide(prs, "Актуальность темы", bullets=[
        "Рестораны без собственного сайта теряют клиентов конкурентам",
        "Заказы по телефону и в мессенджерах — медленно и с ошибками",
        "Агрегаторы доставки берут высокие комиссии",
        "Гость ожидает: меню, цены и заказ — в несколько кликов с телефона",
    ], footer="Ресторан «Анталия» нуждается в собственной цифровой платформе")

    add_content_slide(prs, "Цель и задачи", paragraphs=[
        "Цель: разработать сайт ресторана «Анталия» с автоматизацией онлайн-заказа и бронирования столиков.",
        "Задачи:",
    ], bullets=[
        "Изучить понятие и классификацию сайтов",
        "Проанализировать инструментарий разработки",
        "Описать этапы создания сайта",
        "Провести подготовительные работы для «Анталии»",
        "Спроектировать структуру и интерфейсы",
        "Реализовать и развернуть готовый продукт",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(9), Inches(0.7), "Объект, предмет, база", 28, True, GOLD)
    cards = [
        ("Объект", "Процесс разработки корпоративного сайта для предприятия общественного питания"),
        ("Предмет", "Технология создания сайта для ресторана «Анталия»"),
        ("Проблема", "Противоречие между потребностью гостей в цифровом сервисе и отсутствием полнофункционального сайта"),
        ("База", "Ресторан турецкой кухни «Анталия»"),
    ]
    positions = [(0.6, 1.2), (5.2, 1.2), (0.6, 3.8), (5.2, 3.8)]
    for (x, y), (title, body) in zip(positions, cards):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.2), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD
        shape.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = GOLD
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(13)
        p2.font.color.rgb = WHITE

    add_content_slide(prs, "Проблема AS-IS (до сайта)", bullets=[
        "Заказы — по телефону и в мессенджерах",
        "Приём заказа: 5–10 минут, очередь на линии в часы пик",
        "Ошибки при диктовке адреса и состава заказа",
        "Нет единой базы клиентов и истории заказов",
        "Бронирование — только звонком администратору",
    ], footer="→ Рост нагрузки на персонал, потеря выручки, снижение лояльности")

    add_content_slide(prs, "Решение TO-BE (после внедрения)", bullets=[
        "Гость самостоятельно изучает меню и оформляет заказ онлайн",
        "Данные сразу попадают в MySQL через REST API",
        "Администратор видит заказы в защищённой панели",
        "Онлайн-бронирование столиков + файл календаря .ics",
        "Круглосуточный доступ без ожидания ответа оператора",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(9), Inches(0.7), "Технологический стек", 28, True, GOLD)
    stacks = [
        ("Frontend", "HTML, CSS, JavaScript (ES6)\nBlade, Tailwind CSS 4", RGBColor(0xD9, 0x77, 0x06)),
        ("Backend", "PHP 8.3, Laravel 13\nREST API, Eloquent ORM", RGBColor(0x16, 0xA3, 0x4A)),
        ("БД", "MySQL (production)\nSQLite (разработка)", RGBColor(0x25, 0x63, 0xEB)),
        ("Деплой", "Railway, FrankenPHP\nGitHub, PHPUnit CI", RGBColor(0x93, 0x33, 0xEA)),
    ]
    pos = [(0.6, 1.2), (5.2, 1.2), (0.6, 3.6), (5.2, 3.6)]
    for (x, y), (title, body, accent) in zip(pos, stacks):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.2), Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD
        shape.line.color.rgb = accent
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = accent
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(13)
        p2.font.color.rgb = WHITE
    add_textbox(slide, Inches(0.8), Inches(6.0), Inches(8.6), Inches(0.5), "Архитектура: клиент–сервер, паттерн MVC", 14, False, MUTED)

    add_content_slide(prs, "База данных MySQL", bullets=[
        "menu_categories → menu_items (1:N)",
        "orders → order_items → menu_items",
        "reservations — бронирования столиков",
        "users — администраторы (bcrypt)",
        "Миграции Laravel + сидер с 15 блюдами турецкой кухни",
    ], footer="ER-диаграмма — в приложении к диплому")

    add_content_slide(prs, "Публичная часть сайта", bullets=[
        "Главная — о ресторане, акции, популярные блюда",
        "Меню — каталог с поиском и фильтром по категориям",
        "Корзина и оформление заказа (доставка / самовывоз)",
        "Бронирование столика",
        "Разделы «О нас» и «Контакты»",
    ], footer="Адаптивная вёрстка — приоритет мобильных устройств")

    add_table_slide(prs, "REST API",
        ["Метод", "Маршрут", "Назначение"],
        [
            ("GET", "/api/menu", "Каталог блюд из БД"),
            ("POST", "/api/orders", "Создание заказа"),
            ("POST", "/api/reservations", "Бронирование"),
        ],
        note="Формат обмена — JSON. Валидация на сервере в OrderController")

    add_content_slide(prs, "Корзина и localStorage", bullets=[
        "Нативный JavaScript без React/Vue",
        "Состояние корзины: antalya_cart_v1 в localStorage",
        "Добавление блюд без перезагрузки страницы",
        "Отправка заказа через fetch → POST /api/orders",
        "Данные не теряются при обновлении страницы",
    ])

    add_content_slide(prs, "Админ-панель", bullets=[
        "URL: /admin — доступ только после авторизации",
        "Laravel Auth + сессии + bcrypt",
        "Список заказов, смена статуса (new → accepted → completed)",
        "Управление бронированиями и меню",
        "Защита CSRF, middleware на маршрутах",
    ])

    add_content_slide(prs, "Развёртывание на Railway", bullets=[
        "Облачный хостинг Railway + MySQL",
        "Railpack / FrankenPHP — быстрая обработка запросов",
        "Деплой из GitHub-репозитория",
        "start-container.sh: migrate + db:seed при запуске",
        "HTTPS, переменные APP_KEY, DB_*, MYSQL_*",
    ], footer="Сайт работает в сети Интернет и принимает реальные заказы")

    add_content_slide(prs, "Результаты работы", bullets=[
        "✓ Сайт разработан и развёрнут",
        "✓ Автоматизированы заказ и бронирование",
        "✓ Снижена нагрузка на персонал",
        "✓ Заказы хранятся в единой БД",
        "✓ Настроено автотестирование (PHPUnit, GitHub Actions)",
    ], footer="Практическая значимость: готовый продукт для реального предприятия")

    add_content_slide(prs, "Перспективы развития", bullets=[
        "Онлайн-оплата (эквайринг)",
        "Личный кабинет и программа лояльности",
        "Интеграция со складским учётом",
        "SMS/Telegram-уведомления о статусе заказа",
        "Аналитика продаж и отчёты",
    ])

    add_title_slide(prs, [
        ("Спасибо за внимание!", 36, True, GOLD, PP_ALIGN.CENTER),
        ("Готов ответить на вопросы", 22, False, WHITE, PP_ALIGN.CENTER),
        ("Демонстрация: главная → меню → корзина → админ-панель", 14, False, MUTED, PP_ALIGN.CENTER),
    ])

    prs.save(OUT)
    print(f"Saved: {OUT}")

if __name__ == "__main__":
    main()
